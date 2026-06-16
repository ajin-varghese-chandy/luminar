# scripts/pptx_to_md.py

import re
from pathlib import Path
from pptx import Presentation
import shutil


def extract_images(slide, slide_number, image_dir, prefix):
    markdown = []

    image_count = 1

    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            image = shape.image

            extension = image.ext
            filename = f"{prefix}slide_{slide_number}_image_{image_count}.{extension}"

            image_path = image_dir / filename

            image_path.write_bytes(image.blob)

            img_path = f"{image_dir.name}/{filename}"
            markdown.extend(
                [
                    "",
                    f'<img src="{img_path}" alt="{filename}" width="800">',
                    "",
                ]
            )

            image_count += 1

    return markdown


def convert_pptx_to_md(pptx_path: Path):
    prs = Presentation(pptx_path)

    markdown = [
        f"# {pptx_path.stem}",
        "",
    ]

    md_path = pptx_path.parent / "README.md"

    stem_slug = re.sub(r"[^\w\-]+", "_", pptx_path.stem.strip())

    image_dir = pptx_path.parent / "images"

    image_dir.mkdir(exist_ok=True)

    img_prefix = f"{stem_slug}-"

    for slide_number, slide in enumerate(prs.slides, start=1):
        markdown.extend(
            [
                "---",
                "",
                f"## Slide {slide_number}",
                "",
            ]
        )

        notes = []

        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue

            text = shape.text.strip()

            if not text:
                continue

            # Title
            if slide.shapes.title and shape == slide.shapes.title:
                markdown.extend(
                    [
                        f"# {text}",
                        "",
                    ]
                )
                continue

            # Bullet points
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    content = paragraph.text.strip()

                    if not content:
                        continue

                    indent = "  " * paragraph.level

                    markdown.append(
                        f"{indent}- {content}"
                    )

        # Images
        markdown.extend(
            extract_images(
                slide,
                slide_number,
                image_dir,
                img_prefix,
            )
        )

        # Speaker Notes
        if slide.has_notes_slide:
            for shape in slide.notes_slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if (
                        text
                        and text.lower()
                        != "click to add notes"
                    ):
                        notes.append(text)

        if notes:
            markdown.extend(
                [
                    "",
                    "### Notes",
                    "",
                ]
            )

            for note in notes:
                markdown.extend(
                    [
                        note,
                        "",
                    ]
                )

    md_path.write_text(
        "\n".join(markdown),
        encoding="utf-8",
    )

    print(f"✓ {pptx_path} → {md_path}")


def main():
    root = Path(".")

    pptx_files = list(root.rglob("*.pptx"))

    if not pptx_files:
        print("No PPTX files found.")
        return

    print(
        f"Found {len(pptx_files)} PPTX file(s).\n"
    )

    for pptx_file in pptx_files:
        try:
            convert_pptx_to_md(pptx_file)
        except Exception as e:
            print(
                f"✗ Failed: {pptx_file}"
            )
            print(f"  Error: {e}\n")


if __name__ == "__main__":
    main()