import re
import sys
from pathlib import Path
from pptx import Presentation


def get_slide_count(prs):
    return len(prs.slides)


def get_image_count(prs):
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:
                count += 1
    return count


def get_text_present(prs):
    """Check that slides have some text content (not just blank/placeholder)."""
    non_empty = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                non_empty += 1
                break
    return non_empty


def count_md_slides(md_text):
    return len(re.findall(r"^## Slide \d+", md_text, re.MULTILINE))


def count_md_images(md_text):
    return len(re.findall(r'<img\s+src="images/', md_text))


def get_md_image_paths(md_text, md_dir):
    paths = []
    for m in re.finditer(r'<img\s+src="([^"]+)"', md_text):
        paths.append((md_dir / m.group(1)).resolve())
    return paths


def compute_stem_slug(stem):
    return re.sub(r"[^\w\-]+", "_", stem.strip())


def test_pptx_conversion(pptx_path):
    stem = pptx_path.stem
    md_path = pptx_path.parent / "README.md"
    prs = Presentation(str(pptx_path))
    parent = pptx_path.parent
    md_text = md_path.read_text(encoding="utf-8")

    # 1. MD file exists
    assert md_path.exists(), f"Missing MD file: {md_path}"

    # 2. All slides present
    md_slides = count_md_slides(md_text)
    pptx_slides = get_slide_count(prs)
    assert md_slides == pptx_slides, (
        f"Slide count mismatch in '{stem}': "
        f"MD has {md_slides}, PPTX has {pptx_slides}"
    )

    # 3. All images extracted
    md_imgs = count_md_images(md_text)
    pptx_imgs = get_image_count(prs)
    assert md_imgs == pptx_imgs, (
        f"Image count mismatch in '{stem}': "
        f"MD references {md_imgs}, PPTX has {pptx_imgs}"
    )

    # 4. Referenced image files exist on disk
    stem_slug = compute_stem_slug(stem)
    for img_path in get_md_image_paths(md_text, parent):
        assert img_path.exists(), (
            f"Missing image file referenced in '{stem}': {img_path}"
        )
        assert stem_slug in img_path.name, (
            f"Image '{img_path.name}' missing stem prefix '{stem_slug}'"
        )

    print(f"  ✓ {stem}: {pptx_slides} slides, {pptx_imgs} images")


def main():
    root = Path(".")
    pptx_files = sorted(root.rglob("*.pptx"))

    if not pptx_files:
        print("No PPTX files found.")
        sys.exit(0)

    print(f"Testing {len(pptx_files)} PPTX file(s)...\n")

    failures = 0
    for pptx_path in pptx_files:
        try:
            test_pptx_conversion(pptx_path)
        except AssertionError as e:
            print(f"  ✗ {pptx_path.stem}: {e}")
            failures += 1

    print()
    if failures:
        print(f"❌ {failures} test(s) failed.")
        sys.exit(1)
    else:
        print("✅ All tests passed.")
        sys.exit(0)


if __name__ == "__main__":
    main()
